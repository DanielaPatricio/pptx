from pptx import Presentation
from pptx.util import Inches

prs = Presentation('testprs.pptx')

for slide in prs.slides:
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    for shape in slide.shapes:
         shape.height = prs.slide_height
         shape.width = prs.slide_width
         shape.left = round((prs.slide_width - shape.width) / 2)
         shape.top = round((prs.slide_height - shape.height) / 2)

         print(prs.slide_height - shape.width /2)

prs.save('testprs.pptx')




